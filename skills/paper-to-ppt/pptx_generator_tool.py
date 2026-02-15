"""
PPTX Generator Tool - PPT 生成工具
基于 python-pptx，接收结构化 JSON 描述，生成 .pptx 文件
"""

import sys
import os
sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))

from utils.base_tool import BaseTool
from typing import Dict, Any, List, Optional
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE


# 预设配色方案
THEMES = {
    "academic": {
        "bg": RGBColor(0xFF, 0xFF, 0xFF),
        "title_color": RGBColor(0x1A, 0x1A, 0x2E),
        "subtitle_color": RGBColor(0x4A, 0x4A, 0x6A),
        "text_color": RGBColor(0x33, 0x33, 0x33),
        "accent": RGBColor(0x2B, 0x57, 0x97),
        "accent_light": RGBColor(0xE8, 0xEF, 0xF7),
        "divider": RGBColor(0x2B, 0x57, 0x97),
    },
    "dark": {
        "bg": RGBColor(0x1E, 0x1E, 0x2E),
        "title_color": RGBColor(0xF0, 0xF0, 0xF0),
        "subtitle_color": RGBColor(0xBB, 0xBB, 0xCC),
        "text_color": RGBColor(0xDD, 0xDD, 0xDD),
        "accent": RGBColor(0x6C, 0xA0, 0xDC),
        "accent_light": RGBColor(0x2A, 0x2A, 0x40),
        "divider": RGBColor(0x6C, 0xA0, 0xDC),
    },
    "minimal": {
        "bg": RGBColor(0xFA, 0xFA, 0xFA),
        "title_color": RGBColor(0x22, 0x22, 0x22),
        "subtitle_color": RGBColor(0x66, 0x66, 0x66),
        "text_color": RGBColor(0x44, 0x44, 0x44),
        "accent": RGBColor(0xE0, 0x4B, 0x4B),
        "accent_light": RGBColor(0xFD, 0xED, 0xED),
        "divider": RGBColor(0xE0, 0x4B, 0x4B),
    },
}


class PptxGeneratorTool(BaseTool):
    """PPT 生成工具"""

    def execute(self, params: Dict[str, Any]) -> Dict[str, Any]:
        self.validate_params(params, ['output_path', 'slides'])

        output_path = params['output_path']
        slides_data = params['slides']
        theme_name = params.get('theme', 'academic')

        # 兜底：AI 有时会把 slides 当 JSON 字符串传入，自动反序列化
        if isinstance(slides_data, str):
            import json
            try:
                slides_data = json.loads(slides_data)
            except json.JSONDecodeError as e:
                raise ValueError(f"slides 参数不是合法的 JSON 数组: {e}")
        if not isinstance(slides_data, list):
            raise ValueError(f"slides 参数必须是数组，收到了 {type(slides_data).__name__}")

        self.theme = THEMES.get(theme_name, THEMES['academic'])

        # 确保输出目录存在
        out_dir = os.path.dirname(output_path)
        if out_dir:
            os.makedirs(out_dir, exist_ok=True)

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        for slide_data in slides_data:
            self._add_slide(prs, slide_data)

        prs.save(output_path)

        return {
            'output_path': os.path.abspath(output_path),
            'slide_count': len(slides_data),
            'theme': theme_name,
        }

    # ── 分发器 ──────────────────────────────────────────

    def _add_slide(self, prs: Presentation, data: Dict[str, Any]):
        """根据 layout 类型分发到对应的 slide 构建方法"""
        layout = data.get('layout', 'content')
        builders = {
            'title': self._add_title_slide,
            'content': self._add_content_slide,
            'section_header': self._add_section_header_slide,
            'two_column': self._add_two_column_slide,
            'image_text': self._add_image_text_slide,
        }
        builder = builders.get(layout)
        if not builder:
            raise ValueError(f"不支持的 layout: {layout}，可选: {', '.join(builders.keys())}")
        builder(prs, data)

    def _set_slide_bg(self, slide, color: Optional[RGBColor] = None):
        """设置幻灯片背景色"""
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = color or self.theme['bg']

    # ── 各布局构建方法 ─────────────────────────────────

    def _add_title_slide(self, prs: Presentation, data: Dict[str, Any]):
        """封面页：标题 + 副标题 + 底部装饰线"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
        self._set_slide_bg(slide)

        # 装饰线
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(4.5), Inches(3.8), Inches(4.333), Pt(4)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = self.theme['accent']
        line.line.fill.background()

        # 标题
        title = data.get('title', '')
        txBox = slide.shapes.add_textbox(
            Inches(1.5), Inches(1.8), Inches(10.333), Inches(2.0)
        )
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(40)
        p.font.bold = True
        p.font.color.rgb = self.theme['title_color']
        p.alignment = PP_ALIGN.CENTER

        # 副标题
        subtitle = data.get('subtitle', '')
        if subtitle:
            txBox2 = slide.shapes.add_textbox(
                Inches(2.5), Inches(4.2), Inches(8.333), Inches(1.2)
            )
            tf2 = txBox2.text_frame
            tf2.word_wrap = True
            p2 = tf2.paragraphs[0]
            p2.text = subtitle
            p2.font.size = Pt(20)
            p2.font.color.rgb = self.theme['subtitle_color']
            p2.alignment = PP_ALIGN.CENTER

    def _add_content_slide(self, prs: Presentation, data: Dict[str, Any]):
        """内容页：标题 + 要点列表"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._set_slide_bg(slide)

        # 左侧装饰条
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0), Pt(6), Inches(7.5)
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = self.theme['accent']
        bar.line.fill.background()

        # 标题
        title = data.get('title', '')
        txBox = slide.shapes.add_textbox(
            Inches(0.8), Inches(0.4), Inches(11.5), Inches(0.9)
        )
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = self.theme['title_color']

        # 标题下分隔线
        sep = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.8), Inches(1.35), Inches(2.0), Pt(3)
        )
        sep.fill.solid()
        sep.fill.fore_color.rgb = self.theme['accent']
        sep.line.fill.background()

        # 要点列表
        bullets = data.get('bullets', [])
        if bullets:
            txBox2 = slide.shapes.add_textbox(
                Inches(0.8), Inches(1.7), Inches(11.5), Inches(5.2)
            )
            tf2 = txBox2.text_frame
            tf2.word_wrap = True
            for i, bullet in enumerate(bullets):
                p2 = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
                p2.text = bullet
                p2.font.size = Pt(18)
                p2.font.color.rgb = self.theme['text_color']
                p2.space_after = Pt(10)
                p2.level = 0

    def _add_section_header_slide(self, prs: Presentation, data: Dict[str, Any]):
        """章节分隔页：居中大标题 + 描述"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._set_slide_bg(slide, self.theme['accent_light'])

        # 标题
        title = data.get('title', '')
        txBox = slide.shapes.add_textbox(
            Inches(2.0), Inches(2.5), Inches(9.333), Inches(1.5)
        )
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = self.theme['accent']
        p.alignment = PP_ALIGN.CENTER

        # 描述
        subtitle = data.get('subtitle', '')
        if subtitle:
            txBox2 = slide.shapes.add_textbox(
                Inches(3.0), Inches(4.2), Inches(7.333), Inches(1.0)
            )
            tf2 = txBox2.text_frame
            tf2.word_wrap = True
            p2 = tf2.paragraphs[0]
            p2.text = subtitle
            p2.font.size = Pt(18)
            p2.font.color.rgb = self.theme['subtitle_color']
            p2.alignment = PP_ALIGN.CENTER

    def _add_column(self, slide, left: Emu, col_title: str, bullets: List[str]):
        """在指定位置添加一栏内容（标题 + 要点）"""
        y_offset = Inches(1.7)
        col_width = Inches(5.5)

        if col_title:
            txBox = slide.shapes.add_textbox(left, y_offset, col_width, Inches(0.5))
            tf = txBox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = col_title
            p.font.size = Pt(20)
            p.font.bold = True
            p.font.color.rgb = self.theme['accent']
            y_offset = Inches(2.3)

        if bullets:
            txBox2 = slide.shapes.add_textbox(left, y_offset, col_width, Inches(4.5))
            tf2 = txBox2.text_frame
            tf2.word_wrap = True
            for i, bullet in enumerate(bullets):
                p2 = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
                p2.text = bullet
                p2.font.size = Pt(16)
                p2.font.color.rgb = self.theme['text_color']
                p2.space_after = Pt(8)

    def _add_two_column_slide(self, prs: Presentation, data: Dict[str, Any]):
        """双栏页：标题 + 左右两栏要点"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._set_slide_bg(slide)

        # 左侧装饰条
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0), Pt(6), Inches(7.5)
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = self.theme['accent']
        bar.line.fill.background()

        # 标题
        title = data.get('title', '')
        txBox = slide.shapes.add_textbox(
            Inches(0.8), Inches(0.4), Inches(11.5), Inches(0.9)
        )
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = self.theme['title_color']

        # 分隔线
        sep = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.8), Inches(1.35), Inches(2.0), Pt(3)
        )
        sep.fill.solid()
        sep.fill.fore_color.rgb = self.theme['accent']
        sep.line.fill.background()

        # 左栏
        self._add_column(slide, Inches(0.8),
                         data.get('left_title', ''), data.get('left_bullets', []))
        # 右栏
        self._add_column(slide, Inches(7.0),
                         data.get('right_title', ''), data.get('right_bullets', []))

    def _add_image_text_slide(self, prs: Presentation, data: Dict[str, Any]):
        """图文页：左侧图片 + 右侧文字"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._set_slide_bg(slide)

        # 左侧装饰条
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0), Pt(6), Inches(7.5)
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = self.theme['accent']
        bar.line.fill.background()

        # 标题
        title = data.get('title', '')
        txBox = slide.shapes.add_textbox(
            Inches(0.8), Inches(0.4), Inches(11.5), Inches(0.9)
        )
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = self.theme['title_color']

        # 分隔线
        sep = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.8), Inches(1.35), Inches(2.0), Pt(3)
        )
        sep.fill.solid()
        sep.fill.fore_color.rgb = self.theme['accent']
        sep.line.fill.background()

        # 图片（左侧）
        image_path = data.get('image_path', '')
        if image_path and os.path.isfile(image_path):
            slide.shapes.add_picture(
                image_path,
                Inches(0.8), Inches(1.7),
                Inches(5.5), Inches(5.0)
            )

        # 文字（右侧）
        text = data.get('text', '')
        bullets = data.get('bullets', [])
        txt_left = Inches(7.0)
        txt_width = Inches(5.5)

        if text:
            txBox2 = slide.shapes.add_textbox(
                txt_left, Inches(1.7), txt_width, Inches(5.0)
            )
            tf2 = txBox2.text_frame
            tf2.word_wrap = True
            p2 = tf2.paragraphs[0]
            p2.text = text
            p2.font.size = Pt(16)
            p2.font.color.rgb = self.theme['text_color']
        elif bullets:
            txBox2 = slide.shapes.add_textbox(
                txt_left, Inches(1.7), txt_width, Inches(5.0)
            )
            tf2 = txBox2.text_frame
            tf2.word_wrap = True
            for i, b in enumerate(bullets):
                pb = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
                pb.text = b
                pb.font.size = Pt(16)
                pb.font.color.rgb = self.theme['text_color']
                pb.space_after = Pt(8)


if __name__ == '__main__':
    tool = PptxGeneratorTool()
    tool.run()